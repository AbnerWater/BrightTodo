from __future__ import annotations

from datetime import datetime
from io import BytesIO
from types import SimpleNamespace

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from lifetrace.schemas.todo import TodoStatus
from lifetrace.services.agent_import_service import (
    MAX_IMPORT_FILE_BYTES,
    MAX_IMPORT_FILE_COUNT,
    AgentImportFile,
    AgentImportService,
    ImportTodosError,
)

HTTP_BAD_REQUEST = 400
HTTP_CONTENT_TOO_LARGE = 413
EXPECTED_TEXT_TASK_COUNT = 2
IMAGE_CONFIDENCE = 0.82
DOCUMENT_FALLBACK_CONFIDENCE = 0.55
LLM_TEMPERATURE = 0.2
LLM_MAX_TOKENS = 1800


class FakeUnavailableLlmClient:
    def is_available(self) -> bool:
        return False


class FakeDocumentLlmClient:
    def __init__(self, response_text: str) -> None:
        self.response_text = response_text
        self.messages: list[dict[str, str]] = []

    def is_available(self) -> bool:
        return True

    def chat(
        self,
        messages: list[dict[str, str]],
        temperature: float,
        max_tokens: int,
    ) -> str:
        self.messages = messages
        assert temperature == LLM_TEMPERATURE
        assert max_tokens == LLM_MAX_TOKENS
        return self.response_text


class FakeTodoService:
    def __init__(self) -> None:
        self.created = []

    def list_todos(self, **_kwargs: object):
        return {"total": 0, "todos": []}

    def create_todo(self, data):
        self.created.append(data)
        return SimpleNamespace(id=len(self.created), name=data.name, status=data.status.value)


class FailingAfterFirstTodoService(FakeTodoService):
    def __init__(self) -> None:
        super().__init__()
        self.deleted: list[int] = []

    def create_todo(self, data):
        if self.created:
            raise RuntimeError("模拟创建失败")
        return super().create_todo(data)

    def delete_todo(self, todo_id: int) -> None:
        self.deleted.append(todo_id)


def _reference_time() -> datetime:
    return datetime.fromisoformat("2026-05-19T10:00:00+08:00")


def _service(document_llm_client=None) -> AgentImportService:
    return AgentImportService(
        document_llm_client=document_llm_client or FakeUnavailableLlmClient()
    )


def _txt_file(content: str, name: str = "tasks.txt") -> AgentImportFile:
    return AgentImportFile(
        file_name=name,
        mime_type="text/plain",
        content=content.encode("utf-8"),
    )


def _build_docx_bytes() -> bytes:
    document = Document()
    document.add_paragraph("明天下午三点前完成软件工程报告，大概需要两个小时，紧急")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "后天回复导师邮件"
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _build_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "课程项目"
    sheet.append(["事项", "说明"])
    sheet.append(["明天下午三点前完成 Excel 数据整理", "大概需要两个小时，紧急"])
    sheet.append(["后天准备课程项目展示 PPT", "整理演示材料"])
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _build_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    title_box.text = "明天完成项目汇报材料，大概需要两个小时"
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(6), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = "后天回复导师修改意见"
    table.cell(0, 1).text = "紧急"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_pdf_bytes() -> bytes:
    writer = PdfWriter()
    page = writer.add_blank_page(width=240, height=240)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    resources = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )
    page[NameObject("/Resources")] = resources
    content = DecodedStreamObject()
    content.set_data(b"BT /F1 12 Tf 10 120 Td (TODO finish report tomorrow) Tj ET")
    page[NameObject("/Contents")] = writer._add_object(content)

    buffer = BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def _build_pdf_bytes_with_text(text: str) -> bytes:
    writer = PdfWriter()
    page = writer.add_blank_page(width=360, height=240)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    resources = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )
    page[NameObject("/Resources")] = resources
    content = DecodedStreamObject()
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    content.set_data(f"BT /F1 12 Tf 10 120 Td ({escaped}) Tj ET".encode("ascii"))
    page[NameObject("/Contents")] = writer._add_object(content)

    buffer = BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def test_import_text_file_extracts_multiple_tasks_without_creating() -> None:
    fake_todo_service = FakeTodoService()

    response = _service().import_files(
        files=[
            _txt_file(
                "- 明天下午三点前完成数学作业，大概需要两个小时\n"
                "- 后天回复导师邮件，紧急"
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=fake_todo_service,
    )

    assert response.file_results[0].status == "success"
    assert response.file_results[0].extracted_count == EXPECTED_TEXT_TASK_COUNT
    assert {task.task_title for task in response.extracted_tasks} == {
        "完成数学作业",
        "回复导师邮件",
    }
    assert fake_todo_service.created == []


def test_import_docx_extracts_paragraph_and_table_tasks() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="tasks.docx",
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                content=_build_docx_bytes(),
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    titles = {task.task_title for task in response.extracted_tasks}
    assert "完成软件工程报告" in titles
    assert "回复导师邮件" in titles
    assert "完成软件工程报告" in response.raw_text_preview


def test_import_xlsx_extracts_sheet_rows_as_tasks() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="course.xlsx",
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                content=_build_xlsx_bytes(),
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    titles = {task.task_title for task in response.extracted_tasks}
    assert "完成 Excel 数据整理" in titles
    assert "准备课程项目展示 PPT" in titles
    assert "课程项目" in response.raw_text_preview


def test_import_pptx_extracts_slide_text_and_table_tasks() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="course.pptx",
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                content=_build_pptx_bytes(),
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    titles = {task.task_title for task in response.extracted_tasks}
    assert "完成项目汇报材料" in titles
    assert "回复导师修改意见" in titles
    assert "幻灯片 1" in response.raw_text_preview


def test_import_pdf_extracts_text_task() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="tasks.pdf",
                mime_type="application/pdf",
                content=_build_pdf_bytes(),
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "success"
    assert response.extracted_tasks
    assert response.extracted_tasks[0].source_file == "tasks.pdf"


def test_import_generic_pdf_uses_llm_to_infer_tasks() -> None:
    fake_llm = FakeDocumentLlmClient(
        """
        {
          "todos": [
            {
              "title": "准备课程项目展示",
              "description": "根据项目说明整理展示材料",
              "priority": "medium",
              "due": null,
              "duration": "PT2H",
              "source_text": "Course project requires a group presentation and final report.",
              "confidence": 0.83
            }
          ]
        }
        """
    )

    response = _service(fake_llm).import_files(
        files=[
            AgentImportFile(
                file_name="syllabus.pdf",
                mime_type="application/pdf",
                content=_build_pdf_bytes_with_text(
                    "Course project requires a group presentation and final report."
                ),
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "success"
    assert response.file_results[0].extracted_count == 1
    assert response.extracted_tasks[0].task_title == "准备课程项目展示"
    assert response.extracted_tasks[0].duration == "PT2H"
    assert response.extracted_tasks[0].source_file_index == 0
    assert "Course project requires" in fake_llm.messages[1]["content"]


def test_import_llm_task_gets_initial_duration_when_missing() -> None:
    fake_llm = FakeDocumentLlmClient(
        """
        {
          "todos": [
            {
              "title": "阅读并整理论文要点",
              "description": "整理为课堂讨论笔记",
              "priority": "low",
              "source_text": "Students will discuss the reading in class.",
              "confidence": 0.76
            }
          ]
        }
        """
    )

    response = _service(fake_llm).import_files(
        files=[_txt_file("Students will discuss the reading in class.")],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.extracted_tasks[0].task_title == "阅读并整理论文要点"
    assert response.extracted_tasks[0].duration == "PT1H"
    assert response.extracted_tasks[0].priority == "low"


def test_import_llm_empty_duration_falls_back_to_estimate() -> None:
    fake_llm = FakeDocumentLlmClient(
        """
        {
          "todos": [
            {
              "title": "准备课程项目报告",
              "priority": "medium",
              "duration": "PT",
              "source_text": "Students should prepare a final project report.",
              "confidence": 0.78
            }
          ]
        }
        """
    )

    response = _service(fake_llm).import_files(
        files=[_txt_file("Students should prepare a final project report.")],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.extracted_tasks[0].task_title == "准备课程项目报告"
    assert response.extracted_tasks[0].duration == "PT2H"


def test_import_generic_document_has_local_fallback_task_when_llm_unavailable() -> None:
    response = _service().import_files(
        files=[
            _txt_file(
                "Course project requires a group presentation and final report.",
                name="course-project.txt",
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.extracted_tasks[0].task_title == "准备项目展示和报告"
    assert response.extracted_tasks[0].duration == "PT2H"
    assert response.extracted_tasks[0].confidence == DOCUMENT_FALLBACK_CONFIDENCE


def test_import_corrupted_pdf_returns_file_failure() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="broken.pdf",
                mime_type="application/pdf",
                content=b"not a pdf",
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "failed"
    assert response.file_results[0].error_code == "PDF_PARSE_FAILED"
    assert response.extracted_tasks == []


def test_import_corrupted_docx_returns_file_failure() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="broken.docx",
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                content=b"not a docx",
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "failed"
    assert response.file_results[0].error_code == "DOCX_PARSE_FAILED"
    assert response.extracted_tasks == []


def test_import_corrupted_spreadsheet_returns_file_failure() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="broken.xlsx",
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                content=b"not an excel file",
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "failed"
    assert response.file_results[0].error_code == "SPREADSHEET_PARSE_FAILED"
    assert response.extracted_tasks == []


def test_import_corrupted_presentation_returns_file_failure() -> None:
    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="broken.pptx",
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                content=b"not a pptx file",
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "failed"
    assert response.file_results[0].error_code == "PRESENTATION_PARSE_FAILED"
    assert response.extracted_tasks == []


def test_import_image_uses_existing_vision_extraction(monkeypatch: pytest.MonkeyPatch) -> None:
    class FakeLlmClient:
        def is_available(self) -> bool:
            return True

    monkeypatch.setattr(
        "lifetrace.services.agent_import_service.get_llm_client",
        lambda: FakeLlmClient(),
    )
    monkeypatch.setattr(
        "lifetrace.services.agent_import_service._call_vision_model_with_base64",
        lambda **_kwargs: [
            {
                "title": "提交截图报告",
                "source_text": "明天提交截图报告，紧急",
                "confidence": IMAGE_CONFIDENCE,
            }
        ],
    )

    response = _service().import_files(
        files=[
            AgentImportFile(
                file_name="capture.png",
                mime_type="image/png",
                content=b"\x89PNG\r\n\x1a\n",
            )
        ],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.file_results[0].status == "success"
    assert response.extracted_tasks[0].task_title == "提交截图报告"
    assert response.extracted_tasks[0].confidence == IMAGE_CONFIDENCE


def test_import_create_todos_creates_draft_items() -> None:
    fake_todo_service = FakeTodoService()

    response = _service().import_files(
        files=[_txt_file("明天完成课程展示材料，大概需要两个小时")],
        reference_time=_reference_time(),
        create_todos=True,
        todo_service=fake_todo_service,
    )

    assert response.created_todos[0].name == "完成课程展示材料"
    assert fake_todo_service.created[0].status == TodoStatus.DRAFT
    assert fake_todo_service.created[0].tags == ["文件导入", "AI解析"]


def test_rule_based_import_adds_initial_duration_estimate() -> None:
    response = _service().import_files(
        files=[_txt_file("后天回复导师邮件，紧急")],
        reference_time=_reference_time(),
        create_todos=False,
        todo_service=FakeTodoService(),
    )

    assert response.extracted_tasks[0].task_title == "回复导师邮件"
    assert response.extracted_tasks[0].duration == "PT30M"


def test_import_create_todos_rolls_back_partial_success() -> None:
    fake_todo_service = FailingAfterFirstTodoService()

    with pytest.raises(ImportTodosError) as exc_info:
        _service().import_files(
            files=[_txt_file("- 明天完成课程展示材料\n- 后天回复导师邮件，紧急")],
            reference_time=_reference_time(),
            create_todos=True,
            todo_service=fake_todo_service,
        )

    assert exc_info.value.error_code == "TODO_CREATE_FAILED"
    assert fake_todo_service.deleted == [1]


@pytest.mark.parametrize(
    ("files", "error_code", "status_code"),
    [
        ([], "INVALID_INPUT", HTTP_BAD_REQUEST),
        (
            [
                AgentImportFile(
                    file_name=f"task-{index}.txt",
                    mime_type="text/plain",
                    content=b"todo",
                )
                for index in range(MAX_IMPORT_FILE_COUNT + 1)
            ],
            "TOO_MANY_FILES",
            HTTP_BAD_REQUEST,
        ),
        ([AgentImportFile("empty.txt", "text/plain", b"")], "EMPTY_FILE", HTTP_BAD_REQUEST),
        (
            [
                AgentImportFile(
                    "large.txt",
                    "text/plain",
                    b"x" * (MAX_IMPORT_FILE_BYTES + 1),
                )
            ],
            "FILE_TOO_LARGE",
            HTTP_CONTENT_TOO_LARGE,
        ),
        (
            [AgentImportFile("archive.zip", "application/zip", b"zip")],
            "INVALID_FILE_TYPE",
            HTTP_BAD_REQUEST,
        ),
    ],
)
def test_import_validation_errors(
    files: list[AgentImportFile],
    error_code: str,
    status_code: int,
) -> None:
    with pytest.raises(ImportTodosError) as exc_info:
        _service().import_files(
            files=files,
            reference_time=_reference_time(),
            create_todos=False,
            todo_service=FakeTodoService(),
        )

    assert exc_info.value.error_code == error_code
    assert exc_info.value.status_code == status_code
